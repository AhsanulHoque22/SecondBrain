"""Builds the slide deck for the schedule's 'Prep slides for guest lecture'
task — this represents the planner actually executing that flexible task,
not just naming it.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG = RGBColor(0x1E, 0x20, 0x26)
FG = RGBColor(0xEE, 0xF0, 0xF3)
ACCENT = RGBColor(0x6E, 0xC5, 0xE8)
MUTED = RGBColor(0x9A, 0xA3, 0xAF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def add_text(slide, text, left, top, width, height, size, color=FG, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Hack"
    return box


def add_bullets(slide, items, left, top, width, height, size=22, color=FG):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"›  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Hack"
        p.space_after = Pt(14)


# Slide 1 — Title
s = add_slide()
add_text(s, "AGENTIC AI", 1, 2.3, 11.3, 1.2, 54, ACCENT, bold=True)
add_text(s, "What makes an AI agent ‘agentic’?", 1, 3.4, 11.3, 0.8, 26, FG)
add_text(s, "Guest Lecture  ·  Built live with Claude Code", 1, 6.6, 11.3, 0.5, 16, MUTED)

# Slide 2 — Chatbot vs Agent
s = add_slide()
add_text(s, "From Chatbot to Agent", 1, 0.6, 11, 0.9, 34, ACCENT, bold=True)
add_bullets(s, [
    "Chatbot: one prompt in, one answer out. No memory of whether it worked.",
    "Agent: perceives state, plans a step, acts with real tools, checks the result.",
    "The difference isn’t model size. It’s the loop around the model.",
], 1, 2.0, 11.2, 4.5)

# Slide 3 — The Core Loop
s = add_slide()
add_text(s, "The Core Loop", 1, 0.6, 11, 0.9, 34, ACCENT, bold=True)
add_bullets(s, [
    "1. Plan: break the goal into a concrete next step",
    "2. Act: call a real tool (write a file, run code, hit an API)",
    "3. Observe: read back what actually happened",
    "4. Verify: check the result against the actual requirement",
    "5. Repeat until the requirement is genuinely met",
], 1, 2.0, 11.2, 4.5)

# Slide 4 — Today's live demo
s = add_slide()
add_text(s, "Today’s Live Demo", 1, 0.6, 11, 0.9, 34, ACCENT, bold=True)
add_bullets(s, [
    "Built a daily planner from constraints: prayer times, class, tasks",
    "First draft had a real scheduling conflict. The validator caught it.",
    "Diagnosed the bug, fixed the logic, re-validated until it passed",
    "Pushed the finished plan to Google Calendar with reminders",
    "Generated this exact slide deck as one of the plan’s own tasks",
    "Drafted a reply to a sample email, another task done autonomously",
], 1, 2.0, 11.2, 5.0, size=20)

# Slide 5 — Takeaway
s = add_slide()
add_text(s, "Key Takeaway", 1, 0.6, 11, 0.9, 34, ACCENT, bold=True)
add_text(
    s,
    "Agentic = tool use + iteration + judgment under real constraints.\n"
    "Not a bigger answer. A system that checks its own work.",
    1, 2.3, 11.2, 3, 26, FG,
)

prs.save("guest_lecture_agentic_ai.pptx")
print("saved guest_lecture_agentic_ai.pptx")
