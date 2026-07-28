"""Builds the final teaching presentation for the lecture: title, agentic-AI
concept intro, then every build step with its prompt, screenshot, and an
explanation of what Claude was actually reasoning/planning/acting on.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

SHOTS = os.path.join(os.path.dirname(__file__), "..", "screenshots")

BG = RGBColor(0x1A, 0x1C, 0x22)
FG = RGBColor(0xEE, 0xF0, 0xF3)
ACCENT = RGBColor(0x6E, 0xC5, 0xE8)
GREEN = RGBColor(0x6C, 0xD1, 0x96)
MUTED = RGBColor(0x9A, 0xA3, 0xAF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_text(slide, text, left, top, width, height, size, color=FG, bold=False, italic=False, align=PP_ALIGN.LEFT, font="Hack"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, items, left, top, width, height, size=20, color=FG):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"›  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Hack"
        p.space_after = Pt(10)


def add_image_fit(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w, h = max_w, max_w / ratio
    else:
        h, w = max_h, max_h * ratio
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def step_slide(step_no, title, prompt, screenshot, reasoning, image2=None):
    s = add_slide()
    add_text(s, f"Step {step_no}", 0.7, 0.35, 4, 0.5, 16, ACCENT, bold=True)
    add_text(s, title, 0.7, 0.75, 11.9, 0.7, 28, FG, bold=True)
    add_text(s, f"PROMPT:  “{prompt}”", 0.7, 1.5, 11.9, 0.6, 15, GREEN, italic=True)
    add_text(s, reasoning, 0.7, 2.05, 11.9, 1.0, 14, MUTED)
    if image2:
        add_image_fit(s, screenshot, 0.5, 3.15, 6.0, 4.0)
        add_image_fit(s, image2, 6.8, 3.15, 6.0, 4.0)
    else:
        add_image_fit(s, screenshot, 0.9, 3.15, 11.5, 4.0)
    return s


def shot(name):
    return os.path.join(SHOTS, name)


# Slide 1 — Title
s = add_slide()
add_text(s, "AGENTIC AI, LIVE", 1, 2.5, 11.3, 1.2, 50, ACCENT, bold=True)
add_text(s, "Building a Daily Planner With Claude Code", 1, 3.6, 11.3, 0.7, 24, FG)
add_text(s, "Guest Lecture Demo  ·  Every step below actually ran", 1, 6.6, 11.3, 0.5, 15, MUTED)

# Slide 2 — What is agentic AI
s = add_slide()
add_text(s, "What Is Agentic AI?", 0.7, 0.5, 11.5, 0.9, 32, ACCENT, bold=True)
add_bullets(s, [
    "A chatbot answers once and stops. It has no idea if it was right.",
    "An agent plans a step, acts with real tools, checks the result, and repeats.",
    "Today: watch that loop build something real, end to end, live.",
], 0.9, 2.1, 11.3, 4)

# Slide 3 — The core loop
s = add_slide()
add_text(s, "The Core Loop", 0.7, 0.5, 11.5, 0.9, 32, ACCENT, bold=True)
add_bullets(s, [
    "1. Plan: break the goal into a concrete next step",
    "2. Act: call a real tool (write code, run it, hit an API)",
    "3. Observe: read back what actually happened",
    "4. Verify: check the result against the real requirement",
    "5. Repeat until it is genuinely done, not just attempted",
], 0.9, 2.0, 11.3, 4.5)

# Slide 4 — What we're building
s = add_slide()
add_text(s, "What We're Building", 0.7, 0.5, 11.5, 0.9, 32, ACCENT, bold=True)
add_bullets(s, [
    "A daily planner from real constraints: prayer times, a class, seven flexible tasks",
    "A validator that enforces the constraints, not a script that just prints a plan",
    "Then: push it to Google Calendar, prep the class slides, draft a client reply",
], 0.9, 2.1, 11.3, 4)

# Step slides
step_slide(
    1, "Set the Rules",
    "Write constraints.json with today's fixed blocks and flexible tasks, plus validate.py to check for overlaps and enforce the fixed blocks.",
    shot("01_constraints_setup.png"),
    "Before generating anything, Claude writes down what ‘correct’ means as code. That turns a vague request into something it can actually check its own work against later.",
)

step_slide(
    2, "First Draft",
    "Write schedule.py to generate today's schedule from constraints.json, then run it.",
    shot("02b_schedule_v1_condensed.png"),
    "Claude writes an honest first attempt: place fixed blocks, then slot flexible tasks in file order with no look-ahead. This is what a plain greedy scheduler actually produces, not a staged mistake.",
)

step_slide(
    3, "Observe the Failure",
    "Run validate.py against the generated schedule.",
    shot("03_validate_v1_fail.png"),
    "Claude runs the same check a human reviewer would run. The output is not a final answer, it is a signal: three real overlaps, including one on top of Fajr prayer. This is the ‘observe’ step of the loop.",
)

step_slide(
    4, "Diagnose and Fix",
    "The validator failed. Diagnose why, fix schedule.py, then regenerate the schedule.",
    shot("04b_schedule_v2_condensed.png"),
    "Claude traces every failure to one root cause: no priority order and no look-ahead. The fix sorts tasks by priority and scans forward for a truly free slot, rather than patching each conflict by hand. That is judgment, not guesswork.",
)

step_slide(
    5, "Verify",
    "Re-run the validator against the fixed schedule.",
    shot("05_validate_v2_pass.png"),
    "Only now does Claude call it done, because the same objective check that failed before now passes. The loop closes on evidence, not on confidence.",
)

step_slide(
    6, "Act in the Real World",
    "Push today's validated schedule to Google Calendar with reminders.",
    shot("06_calendar_morning.jpg"),
    "Agentic does not stop at generating text. Claude calls a real calendar API, so the plan becomes something that will actually remind him at the right time.",
    image2=shot("07_calendar_afternoon.jpg"),
)

step_slide(
    7, "Execute a Task From Its Own Plan",
    "One of today's tasks is prepping slides for the guest lecture. Generate that slide deck.",
    shot("08_class_deck_title.png"),
    "The plan is not a list Claude wrote and handed off. Claude carries out one of its own line items right there, as a real, checkable deliverable.",
    image2=shot("09_class_deck_recap.png"),
)

step_slide(
    8, "Draft the Reply",
    "Another task is replying to a client email. Draft the reply.",
    shot("10_mock_email_task.png"),
    "Same pattern again: read the actual request, act on it, and produce something the user only has to review, not write from scratch.",
)

# Final slide — takeaway
s = add_slide()
add_text(s, "Key Takeaway", 0.7, 0.6, 11.5, 0.9, 34, ACCENT, bold=True)
add_bullets(s, [
    "Agentic = tool use + iteration + judgment under real constraints.",
    "Every step here was checked against something real: a validator, a calendar API, an actual task in the plan.",
    "Try it yourselves: give Claude Code a goal with a checkable definition of done, and watch what it does with the loop.",
], 0.9, 2.2, 11.3, 4, size=22)

out_path = os.path.join(os.path.dirname(__file__), "Agentic_AI_Live_Demo.pptx")
prs.save(out_path)
print(f"saved {out_path}")
