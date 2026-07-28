"""Builds a slide deck from the n8n beginner tutorial article, with every
screenshot from the article and a step-by-step walkthrough.

Source: "n8n Tutorial for Beginners: Build Your First No-Code Automation
Workflow Step by Step" by Rizwanhoda, Towards AI.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

HERE = os.path.dirname(__file__)
IMG = lambda name: os.path.join(HERE, "images", name)

BG = RGBColor(0x1A, 0x1C, 0x22)
FG = RGBColor(0xEE, 0xF0, 0xF3)
ACCENT = RGBColor(0xF2, 0x6B, 0x5A)  # n8n-ish coral
MUTED = RGBColor(0x9A, 0xA3, 0xAF)
GREEN = RGBColor(0x6C, 0xD1, 0x96)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_text(slide, text, left, top, width, height, size, color=FG, bold=False, italic=False, align=PP_ALIGN.LEFT, font="Hack"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, items, left, top, width, height, size=20, color=FG):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"›  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Hack"
        p.space_after = Pt(10)


def add_image_fit(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w, h = max_w, max_w / ratio
    else:
        h, w = max_h, max_h * ratio
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def kicker_title(slide, kicker, title):
    if kicker:
        add_text(slide, kicker, 0.7, 0.35, 10, 0.5, 16, ACCENT, bold=True)
        add_text(slide, title, 0.7, 0.75, 11.9, 0.8, 28, FG, bold=True)
    else:
        add_text(slide, title, 0.7, 0.5, 11.9, 0.9, 32, ACCENT, bold=True)


def step_slide(kicker, title, note, image, image2=None):
    s = add_slide()
    kicker_title(s, kicker, title)
    top_text = 1.55 if kicker else 1.5
    add_text(s, note, 0.7, top_text, 11.9, 1.0, 15, MUTED)
    img_top = top_text + 0.85
    img_h = 7.5 - img_top - 0.35
    if image2:
        add_image_fit(s, image, 0.5, img_top, 6.0, img_h)
        add_image_fit(s, image2, 6.8, img_top, 6.0, img_h)
    else:
        add_image_fit(s, image, 1.2, img_top, 10.9, img_h)
    return s


# 1 — Title
s = add_slide()
add_image_fit(s, IMG("06_hero.jpg"), 7.6, 0, 5.73, 7.5)
add_text(s, "N8N TUTORIAL", 0.8, 2.5, 6.5, 1.0, 44, ACCENT, bold=True)
add_text(s, "Build Your First No-Code\nAutomation Workflow,\nStep by Step", 0.8, 3.4, 6.5, 2.2, 26, FG)
add_text(s, "Based on the article by Rizwanhoda for Towards AI", 0.8, 6.7, 6.5, 0.5, 13, MUTED)

# 2 — What is n8n
s = add_slide()
kicker_title(s, None, "What Is n8n?")
add_bullets(s, [
    "A visual, no-code workflow automation platform (pronounced ‘n-eight-n’)",
    "Connect apps and services by linking nodes on a canvas, no code required",
    "Free and open-source when you self-host it on your own machine",
], 0.9, 2.0, 11.3, 4)

# 3 — Why it matters
s = add_slide()
kicker_title(s, None, "Why It Matters for Beginners")
add_bullets(s, [
    "No coding background needed: drag, drop, and connect nodes",
    "Scales from a two-step automation to a full backend workflow",
    "Tasks that used to take hours by hand now run in one click",
], 0.9, 2.0, 11.3, 4)

# 4 — Installing, two paths
s = add_slide()
kicker_title(s, None, "Installing n8n: Two Paths")
add_bullets(s, [
    "Local (self-hosted): needs Node.js first, then full control and privacy",
    "n8n Cloud: sign up and start building in minutes, zero setup",
    "This tutorial uses the cloud version for every screenshot ahead",
], 0.9, 1.7, 11.3, 2.1, size=18)
add_text(s, "$ npm install -g n8n\n$ n8n start\n# then open http://localhost:5678\n\n$ npm update -g n8n   # to upgrade later", 0.9, 4.0, 11.3, 2.5, 18, GREEN)

# 5 — The cloud editor
step_slide(
    None, "The Cloud Editor",
    "After signing up and answering a few setup questions, click “Start Automating”. You land on a blank canvas with one node already placed, ready to build.",
    IMG("32_cloud_editor.png"),
)

# 6 — Agenda
s = add_slide()
kicker_title(s, None, "What We're Building")
add_bullets(s, [
    "1. Form Trigger  →  2. Google Sheets  →  3. Switch",
    "4. Gmail  →  5. Duplicate Gmail  →  6. Merge  →  7. Test",
    "Form submission → data saved → the right email sent, automatically",
], 0.9, 2.0, 11.3, 4)

# Step 1 — Form Trigger (5 slides)
step_slide(
    "Step 1a", "Start the Workflow",
    "Click “Start from Scratch”, then “Add First Step” to open the node search panel and begin the flow.",
    IMG("38_add_first_step.png"),
)
step_slide(
    "Step 1b", "Add the Form Trigger",
    "Search “On Submission of Form” in the panel and select it as the trigger for this workflow.",
    IMG("40_form_trigger_selected.png"),
)
step_slide(
    "Step 1c", "Define the Form Fields",
    "Add the fields to collect from the user, for example Full Name, Industry, Qualification, and Job Title.",
    IMG("42_form_fields.png"),
)
step_slide(
    "Step 1d", "Run It, See the Form",
    "Executing this step generates a live, working form based on the fields you just defined.",
    IMG("44_generated_form.png"),
)
step_slide(
    "Step 1e", "Inspect the Output",
    "Submitted data appears as JSON in the output panel. Switch to Table or Schema view to explore it differently.",
    IMG("46_json_output.png"),
)

# Step 2 — Google Sheets (3 slides)
step_slide(
    "Step 2a", "Add a Google Sheets Node",
    "Search “Google Sheets” and choose the “Append Row” action, so every submission becomes a new row.",
    IMG("50_sheets_append_row.png"),
)
step_slide(
    "Step 2b", "Point It at a Sheet",
    "Set Resource to Document, pick the target spreadsheet, and connect your Google account under Credentials.",
    IMG("54_sheets_config.png"),
)
step_slide(
    "Step 2c", "Map the Columns",
    "Set Resource to Sheet, then drag each form field onto its matching column before executing.",
    IMG("57_sheets_map_values.png"),
)

# Step 3 — Switch
step_slide(
    "Step 3", "Add a Switch Node",
    "Add a Switch node and set routing rules so the flow branches based on a field value, like the selected option.",
    IMG("61_switch_node.png"),
)

# Step 4 — Gmail (4 slides)
step_slide(
    "Step 4a", "Add Gmail, Connect the Account",
    "Search “Gmail”, choose “Send Message”, and authorize it via the pencil icon under Credentials.",
    IMG("64_gmail_credentials.png"),
)
step_slide(
    "Step 4b", "Write the Message",
    "Set Resource to Message, enter the recipient's address, and write the email content to send.",
    IMG("66_gmail_message_config.png"),
)
step_slide(
    "Step 4c", "Send It",
    "Execute the node and check the inbox. The email sent through the automation should be waiting there.",
    IMG("68_gmail_execute.png"),
)
step_slide(
    "Step 4d", "Configure the Second Branch",
    "Return to the other Gmail node the Switch can route to, and fill in its message details as well.",
    IMG("70_gmail_second_node.png"),
)

# Step 5 — Duplicate
step_slide(
    "Step 5", "Duplicate for a Second Response",
    "Right-click the Gmail node, choose Duplicate, and change the message text to reflect the other outcome.",
    IMG("74_duplicate_gmail.png"),
)

# Step 6 — Merge (2 slides)
step_slide(
    "Step 6a", "Add a Merge Node",
    "Click the + icon on a Gmail node, search “Merge”, and add it to combine both branches into one path.",
    IMG("78_merge_node_config.png"),
)
step_slide(
    "Step 6b", "Connect Both Branches",
    "Drag both Gmail outputs into the Merge node. Form submission, sheet update, and the right email are now one flow.",
    IMG("81_merge_connect_branches.png"),
)

# Step 7 — Test
step_slide(
    "Step 7", "Test End to End",
    "Resubmit the form and confirm two things: the row lands in Google Sheets, and the correct email arrives in the inbox.",
    IMG("87_test_workflow.png"),
)

# Conclusion
s = add_slide()
kicker_title(s, None, "Conclusion")
add_bullets(s, [
    "One trigger, seven nodes, zero code: a complete, working automation",
    "The same pattern (trigger, act, branch, notify) fits almost any repetitive task",
    "Next: try automating something from your own daily work",
], 0.9, 2.0, 11.3, 4)

# Source
s = add_slide()
kicker_title(s, None, "Source")
add_text(
    s,
    "n8n Tutorial for Beginners: Build Your First No-Code\n"
    "Automation Workflow Step by Step\n\n"
    "by Rizwanhoda, published on Towards AI (Medium)",
    0.9, 2.3, 11.3, 2.5, 20, FG,
)
add_text(s, "pub.towardsai.net", 0.9, 5.2, 11.3, 0.5, 15, MUTED)

out = os.path.join(HERE, "n8n_Tutorial_Slides.pptx")
prs.save(out)
print(f"saved {out}, {len(prs.slides)} slides")
